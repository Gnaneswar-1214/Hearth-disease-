from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
import pptx.oxml.ns as nsmap
from lxml import etree

# ── Palette (from the dark ECG theme) ──────────────────────────
BG_DARK   = RGBColor(0x0D, 0x0D, 0x0D)
BG_CARD   = RGBColor(0x1A, 0x08, 0x08)
RED_NEON  = RGBColor(0xFF, 0x2D, 0x2D)
RED_MID   = RGBColor(0xE7, 0x4C, 0x3C)
RED_DARK  = RGBColor(0xC0, 0x39, 0x2B)
WHITE     = RGBColor(0xFF, 0xFF, 0xFF)
OFF_WHITE = RGBColor(0xF0, 0xF0, 0xF0)
GREY      = RGBColor(0xA0, 0xA0, 0xA0)
GREEN     = RGBColor(0x00, 0xE6, 0x76)
AMBER     = RGBColor(0xFF, 0xAB, 0x00)

W = Inches(13.33)   # widescreen width
H = Inches(7.5)     # widescreen height

prs = Presentation()
prs.slide_width  = W
prs.slide_height = H

blank_layout = prs.slide_layouts[6]  # completely blank


# ── Helpers ────────────────────────────────────────────────────

def add_rect(slide, x, y, w, h, fill=None, line=None, line_w=Pt(0)):
    shape = slide.shapes.add_shape(1, x, y, w, h)  # MSO_SHAPE_TYPE.RECTANGLE
    shape.line.width = line_w
    if fill:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill
    else:
        shape.fill.background()
    if line:
        shape.line.color.rgb = line
    else:
        shape.line.fill.background()
    return shape

def add_text(slide, text, x, y, w, h, size=Pt(18), bold=False, color=WHITE,
             align=PP_ALIGN.LEFT, wrap=True):
    txb = slide.shapes.add_textbox(x, y, w, h)
    tf  = txb.text_frame
    tf.word_wrap = wrap
    p   = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size  = size
    run.font.bold  = bold
    run.font.color.rgb = color
    run.font.name  = "Segoe UI"
    return txb

def bg(slide):
    """Full dark background."""
    add_rect(slide, 0, 0, W, H, fill=BG_DARK)

def top_bar(slide, title, subtitle=None):
    """Red accent bar + title at top."""
    add_rect(slide, 0, 0, W, Inches(0.06), fill=RED_NEON)
    add_text(slide, title,
             Inches(0.5), Inches(0.18), Inches(12), Inches(0.7),
             size=Pt(32), bold=True, color=WHITE)
    if subtitle:
        add_text(slide, subtitle,
                 Inches(0.5), Inches(0.88), Inches(12), Inches(0.4),
                 size=Pt(14), color=GREY)

def card(slide, x, y, w, h):
    """Dark card with red border."""
    add_rect(slide, x, y, w, h, fill=BG_CARD, line=RED_NEON, line_w=Pt(1.2))

def bullet_list(slide, items, x, y, w, h, size=Pt(14), color=OFF_WHITE, icon="▸ "):
    txb = slide.shapes.add_textbox(x, y, w, h)
    tf  = txb.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_before = Pt(4)
        run = p.add_run()
        run.text = icon + item
        run.font.size  = size
        run.font.color.rgb = color
        run.font.name  = "Segoe UI"

def red_label(slide, text, x, y, w=Inches(2.5), h=Inches(0.35)):
    add_rect(slide, x, y, w, h, fill=RGBColor(0x3A,0x08,0x08), line=RED_NEON, line_w=Pt(0.8))
    add_text(slide, text, x+Inches(0.1), y+Inches(0.02), w-Inches(0.2), h,
             size=Pt(11), bold=True, color=RED_NEON, align=PP_ALIGN.CENTER)


# ══════════════════════════════════════════════════════════════
# SLIDE 1 — Title
# ══════════════════════════════════════════════════════════════
s = prs.slides.add_slide(blank_layout)
bg(s)
add_rect(s, 0, 0, W, Inches(0.08), fill=RED_NEON)
add_rect(s, 0, H-Inches(0.08), W, Inches(0.08), fill=RED_NEON)

# Centre glow circle
add_rect(s, Inches(4.5), Inches(1.2), Inches(4.3), Inches(4.3),
         fill=RGBColor(0x25,0x05,0x05))

add_text(s, "❤️", Inches(5.8), Inches(1.4), Inches(2), Inches(1.5),
         size=Pt(72), align=PP_ALIGN.CENTER)

add_text(s, "Heart Disease Prediction System",
         Inches(1), Inches(3.1), Inches(11.3), Inches(1.1),
         size=Pt(40), bold=True, color=WHITE, align=PP_ALIGN.CENTER)

add_text(s, "AI-Powered Cardiovascular Risk Assessment",
         Inches(1), Inches(4.15), Inches(11.3), Inches(0.6),
         size=Pt(20), color=RED_NEON, align=PP_ALIGN.CENTER)

add_text(s, "Flask  ·  Scikit-learn  ·  Random Forest  ·  98.5% Accuracy",
         Inches(1), Inches(4.85), Inches(11.3), Inches(0.5),
         size=Pt(14), color=GREY, align=PP_ALIGN.CENTER)

add_text(s, "Built with Python · Deployed on Render",
         Inches(1), Inches(6.8), Inches(11.3), Inches(0.4),
         size=Pt(12), color=RGBColor(0x60,0x60,0x60), align=PP_ALIGN.CENTER)


# ══════════════════════════════════════════════════════════════
# SLIDE 2 — Problem Statement & Objective
# ══════════════════════════════════════════════════════════════
s = prs.slides.add_slide(blank_layout)
bg(s)
top_bar(s, "Problem Statement & Objective",
        "Why early heart disease detection matters")

card(s, Inches(0.4), Inches(1.4), Inches(5.9), Inches(5.5))
add_text(s, "🔴  The Problem",
         Inches(0.6), Inches(1.55), Inches(5.5), Inches(0.45),
         size=Pt(16), bold=True, color=RED_NEON)
bullet_list(s, [
    "Heart disease is the #1 cause of death globally",
    "Late diagnosis leads to poor outcomes",
    "Manual risk assessment is time-consuming",
    "Patients lack easy access to early screening tools",
    "Doctors need decision-support tools for faster triage",
], Inches(0.6), Inches(2.05), Inches(5.5), Inches(4.5), size=Pt(13))

card(s, Inches(6.7), Inches(1.4), Inches(6.2), Inches(5.5))
add_text(s, "🎯  Our Objective",
         Inches(6.9), Inches(1.55), Inches(5.8), Inches(0.45),
         size=Pt(16), bold=True, color=GREEN)
bullet_list(s, [
    "Build a web app for instant heart disease risk prediction",
    "Use ML trained on real clinical data (Kaggle dataset)",
    "Accept 13 standard medical features as input",
    "Return risk probability + clear High/Low label",
    "Provide personalised diet recommendations by age",
    "Deploy publicly so anyone can access it",
], Inches(6.9), Inches(2.05), Inches(5.8), Inches(4.5), size=Pt(13), color=OFF_WHITE)


# ══════════════════════════════════════════════════════════════
# SLIDE 3 — Tech Stack & Architecture
# ══════════════════════════════════════════════════════════════
s = prs.slides.add_slide(blank_layout)
bg(s)
top_bar(s, "Tech Stack & System Architecture")

# 4 tech cards
techs = [
    ("🐍  Python / Flask", "Backend web framework\nRoutes: GET / and POST /predict\nJinja2 server-side rendering", RED_MID),
    ("🤖  Scikit-learn", "Random Forest Classifier\nn_estimators=100, random_state=42\nStandardScaler preprocessing", RGBColor(0x00,0x99,0xFF)),
    ("🌐  HTML / CSS / JS", "Responsive dark ECG theme\n2-column form grid layout\nClient-side validation", GREEN),
    ("☁️  Render (Cloud)", "Free-tier deployment\ngunicorn WSGI server\nAuto-deploy from GitHub", AMBER),
]

for i, (title, body, col) in enumerate(techs):
    cx = Inches(0.35 + i * 3.2)
    card(s, cx, Inches(1.4), Inches(3.0), Inches(2.8))
    add_rect(s, cx, Inches(1.4), Inches(3.0), Inches(0.06), fill=col)
    add_text(s, title, cx+Inches(0.15), Inches(1.55), Inches(2.7), Inches(0.5),
             size=Pt(13), bold=True, color=col)
    add_text(s, body, cx+Inches(0.15), Inches(2.1), Inches(2.7), Inches(1.9),
             size=Pt(11), color=GREY)

# Architecture flow
add_text(s, "System Flow",
         Inches(0.4), Inches(4.4), Inches(4), Inches(0.4),
         size=Pt(14), bold=True, color=RED_NEON)

flow = ["User Browser", "Flask API (app.py)", "StandardScaler", "Random Forest Model", "Result + Diet Page"]
colors_f = [WHITE, RED_NEON, AMBER, GREEN, WHITE]
for i, (label, col) in enumerate(zip(flow, colors_f)):
    fx = Inches(0.35 + i * 2.55)
    add_rect(s, fx, Inches(4.9), Inches(2.3), Inches(0.55),
             fill=RGBColor(0x25,0x08,0x08), line=col, line_w=Pt(1))
    add_text(s, label, fx+Inches(0.05), Inches(4.92), Inches(2.2), Inches(0.5),
             size=Pt(10), bold=True, color=col, align=PP_ALIGN.CENTER)
    if i < 4:
        add_text(s, "→", Inches(2.5 + i*2.55), Inches(4.9), Inches(0.3), Inches(0.55),
                 size=Pt(18), color=RED_NEON, align=PP_ALIGN.CENTER)

add_text(s, "train_model.py  →  model.pkl + scaler.pkl  (offline training, committed to repo)",
         Inches(0.4), Inches(5.7), Inches(12.5), Inches(0.4),
         size=Pt(11), color=GREY)


# ══════════════════════════════════════════════════════════════
# SLIDE 4 — Dataset & Features
# ══════════════════════════════════════════════════════════════
s = prs.slides.add_slide(blank_layout)
bg(s)
top_bar(s, "Dataset & Input Features",
        "Kaggle Heart Disease Dataset · 1025 records · 14 columns")

features = [
    ("age",      "Patient age in years",              "29 – 77"),
    ("sex",      "Biological sex (0=F, 1=M)",         "0 or 1"),
    ("cp",       "Chest pain type",                   "0 – 3"),
    ("trestbps", "Resting blood pressure (mm Hg)",    "94 – 200"),
    ("chol",     "Serum cholesterol (mg/dl)",         "126 – 564"),
    ("fbs",      "Fasting blood sugar > 120 mg/dl",   "0 or 1"),
    ("restecg",  "Resting ECG results",               "0 – 2"),
    ("thalach",  "Max heart rate achieved (bpm)",     "71 – 202"),
    ("exang",    "Exercise induced angina",           "0 or 1"),
    ("oldpeak",  "ST depression (exercise vs rest)",  "0 – 6.2"),
    ("slope",    "Slope of peak exercise ST segment", "0 – 2"),
    ("ca",       "Major vessels (fluoroscopy)",       "0 – 4"),
    ("thal",     "Thalassemia type",                  "0 – 3"),
]

# Header
hx = Inches(0.35)
add_rect(s, hx, Inches(1.35), Inches(12.6), Inches(0.38), fill=RGBColor(0x3A,0x05,0x05))
for txt, cx, cw in [("Feature", Inches(0.5), Inches(2.0)),
                     ("Description", Inches(2.7), Inches(6.5)),
                     ("Range", Inches(9.4), Inches(3.2))]:
    add_text(s, txt, cx, Inches(1.38), cw, Inches(0.35),
             size=Pt(11), bold=True, color=RED_NEON)

for i, (feat, desc, rng) in enumerate(features):
    ry = Inches(1.73 + i * 0.37)
    row_col = RGBColor(0x18,0x06,0x06) if i % 2 == 0 else RGBColor(0x12,0x04,0x04)
    add_rect(s, hx, ry, Inches(12.6), Inches(0.36), fill=row_col)
    add_text(s, feat,  Inches(0.5),  ry+Inches(0.04), Inches(2.0), Inches(0.32), size=Pt(11), bold=True,  color=RED_NEON)
    add_text(s, desc,  Inches(2.7),  ry+Inches(0.04), Inches(6.5), Inches(0.32), size=Pt(11),             color=OFF_WHITE)
    add_text(s, rng,   Inches(9.4),  ry+Inches(0.04), Inches(3.2), Inches(0.32), size=Pt(11),             color=GREY)

add_text(s, "Target: 0 = No Heart Disease  |  1 = Heart Disease Detected",
         Inches(0.35), Inches(6.65), Inches(12.6), Inches(0.4),
         size=Pt(12), bold=True, color=AMBER, align=PP_ALIGN.CENTER)


# ══════════════════════════════════════════════════════════════
# SLIDE 5 — ML Model & Results
# ══════════════════════════════════════════════════════════════
s = prs.slides.add_slide(blank_layout)
bg(s)
top_bar(s, "Machine Learning Model & Results")

# Left — model details
card(s, Inches(0.35), Inches(1.35), Inches(5.8), Inches(5.6))
add_text(s, "🤖  Model Configuration",
         Inches(0.55), Inches(1.5), Inches(5.4), Inches(0.45),
         size=Pt(15), bold=True, color=RED_NEON)
bullet_list(s, [
    "Algorithm: Random Forest Classifier",
    "n_estimators = 100 trees",
    "random_state = 42 (reproducible)",
    "Train / Test split: 80% / 20%",
    "Train size: 820 samples",
    "Test size:  205 samples",
    "Preprocessing: StandardScaler (zero mean, unit variance)",
    "Missing values: rows with NaN dropped",
], Inches(0.55), Inches(2.0), Inches(5.4), Inches(4.5), size=Pt(13))

# Right — metrics
card(s, Inches(6.55), Inches(1.35), Inches(6.4), Inches(2.6))
add_text(s, "📊  Evaluation Metrics",
         Inches(6.75), Inches(1.5), Inches(6.0), Inches(0.45),
         size=Pt(15), bold=True, color=GREEN)

metrics = [("Test Accuracy", "98.54%", GREEN),
           ("Train Size",    "820",    WHITE),
           ("Test Size",     "205",    WHITE)]
for i, (label, val, col) in enumerate(metrics):
    mx = Inches(6.75 + i * 2.05)
    add_rect(s, mx, Inches(2.05), Inches(1.85), Inches(1.6),
             fill=RGBColor(0x05,0x20,0x10) if col==GREEN else RGBColor(0x20,0x08,0x08),
             line=col, line_w=Pt(0.8))
    add_text(s, val,   mx+Inches(0.1), Inches(2.15), Inches(1.65), Inches(0.7),
             size=Pt(26), bold=True, color=col, align=PP_ALIGN.CENTER)
    add_text(s, label, mx+Inches(0.1), Inches(2.85), Inches(1.65), Inches(0.4),
             size=Pt(10), color=GREY, align=PP_ALIGN.CENTER)

# Why Random Forest
card(s, Inches(6.55), Inches(4.1), Inches(6.4), Inches(2.85))
add_text(s, "✅  Why Random Forest?",
         Inches(6.75), Inches(4.25), Inches(6.0), Inches(0.45),
         size=Pt(14), bold=True, color=AMBER)
bullet_list(s, [
    "Handles non-linear relationships well",
    "Robust to outliers and noisy data",
    "Built-in feature importance ranking",
    "No need for feature selection upfront",
    "Excellent accuracy on tabular medical data",
], Inches(6.75), Inches(4.75), Inches(6.0), Inches(2.0), size=Pt(12), color=OFF_WHITE)


# ══════════════════════════════════════════════════════════════
# SLIDE 6 — Web Application Features
# ══════════════════════════════════════════════════════════════
s = prs.slides.add_slide(blank_layout)
bg(s)
top_bar(s, "Web Application Features",
        "Full-stack Flask app with responsive dark ECG UI")

features_app = [
    ("🏠  Home Page",
     ["Project title & description", "Model accuracy stats chips",
      "Animated heartbeat icon", "Start Prediction CTA button"], RED_NEON),
    ("📋  Medical Form",
     ["13 clinical input fields", "3 grouped sections (Personal / Cardiac / Labs)",
      "Dropdowns with human-readable options", "Client-side validation + error messages"], AMBER),
    ("📊  Result Page",
     ["Large bold risk probability %", "3-tier progress bar (Red/Amber/Green)",
      "Heart Disease Detected / No Heart Disease label", "Contextual advice message"], GREEN),
    ("🥗  Diet Plan Panel",
     ["Age-based diet recommendations (3 groups)", "Eat More vs Avoid two-column layout",
      "Extra strict rules for high-risk patients", "Toggle show/hide with smooth animation"], RGBColor(0x00,0xBF,0xFF)),
    ("🎨  UI / UX",
     ["Dark charcoal background (#0d0d0d)", "CSS anatomical heart + animated ECG line",
      "Neon red glow accents throughout", "Fully responsive (mobile-friendly)"], RGBColor(0xFF,0x80,0xFF)),
    ("🔒  Validation & Safety",
     ["All 13 fields required before submit", "Non-numeric input rejected with error",
      "Model file missing → HTTP 500 error page", "Data never stored — local processing only"], RGBColor(0x80,0xFF,0x80)),
]

for i, (title, pts, col) in enumerate(features_app):
    col_i = i % 3
    row_i = i // 3
    cx = Inches(0.3 + col_i * 4.35)
    cy = Inches(1.35 + row_i * 2.95)
    card(s, cx, cy, Inches(4.1), Inches(2.75))
    add_rect(s, cx, cy, Inches(4.1), Inches(0.06), fill=col)
    add_text(s, title, cx+Inches(0.15), cy+Inches(0.12), Inches(3.8), Inches(0.45),
             size=Pt(13), bold=True, color=col)
    bullet_list(s, pts, cx+Inches(0.15), cy+Inches(0.62), Inches(3.8), Inches(2.0),
                size=Pt(11), color=GREY, icon="• ")


# ══════════════════════════════════════════════════════════════
# SLIDE 7 — Output Screenshots (described as mock UI)
# ══════════════════════════════════════════════════════════════
s = prs.slides.add_slide(blank_layout)
bg(s)
top_bar(s, "Application Output — UI Walkthrough")

screens = [
    ("🏠 Home Page", ["Heart Disease Prediction System", "98.5% Accuracy chip", "Animated ❤️ icon", "Start Prediction button"], RED_NEON),
    ("📋 Prediction Form", ["3 sections: Personal / Cardiac / Labs", "Dropdowns + number inputs", "2-column grid layout", "Analyze & Predict button"], AMBER),
    ("📊 Result — High Risk", ["Risk Probability: 78%", "Red glowing progress bar", "Heart Disease Detected (red)", "View Suggested Diet Plan button"], RGBColor(0xFF,0x60,0x60)),
    ("✅ Result — Low Risk", ["Risk Probability: 22%", "Green glowing progress bar", "No Heart Disease (green)", "Low Risk label"], GREEN),
    ("🥗 Diet Plan (High Risk)", ["Strict cardiac diet rules", "Blueberries, omega-3 fish", "Avoid: salt <1500mg/day", "Consult dietitian note"], RGBColor(0x00,0xBF,0xFF)),
    ("🥗 Diet Plan (Age 55+)", ["Soft-cooked vegetables", "Low-sodium soups & stews", "Avoid: processed meats", "Portion control advice"], RGBColor(0xFF,0xAB,0x00)),
]

for i, (title, pts, col) in enumerate(screens):
    col_i = i % 3
    row_i = i // 3
    cx = Inches(0.3 + col_i * 4.35)
    cy = Inches(1.35 + row_i * 2.95)
    # Mock screen frame
    add_rect(s, cx, cy, Inches(4.1), Inches(2.75),
             fill=RGBColor(0x18,0x06,0x06), line=col, line_w=Pt(1.2))
    # Top bar mock
    add_rect(s, cx, cy, Inches(4.1), Inches(0.32), fill=RGBColor(0x2A,0x08,0x08))
    add_rect(s, cx, cy, Inches(4.1), Inches(0.04), fill=col)
    add_text(s, title, cx+Inches(0.12), cy+Inches(0.06), Inches(3.85), Inches(0.28),
             size=Pt(11), bold=True, color=col)
    bullet_list(s, pts, cx+Inches(0.15), cy+Inches(0.45), Inches(3.8), Inches(2.2),
                size=Pt(11), color=GREY, icon="  ✦ ")


# ══════════════════════════════════════════════════════════════
# SLIDE 8 — Conclusion & Future Scope
# ══════════════════════════════════════════════════════════════
s = prs.slides.add_slide(blank_layout)
bg(s)
top_bar(s, "Conclusion & Future Scope")

card(s, Inches(0.35), Inches(1.35), Inches(5.9), Inches(5.6))
add_text(s, "✅  What We Achieved",
         Inches(0.55), Inches(1.5), Inches(5.5), Inches(0.45),
         size=Pt(15), bold=True, color=GREEN)
bullet_list(s, [
    "98.54% accuracy with Random Forest",
    "Full-stack web app (Flask + HTML/CSS/JS)",
    "13-feature clinical input form with validation",
    "Real-time risk probability with visual indicator",
    "Age-based personalised diet recommendations",
    "Deployed live on Render (cloud)",
    "Clean dark ECG-themed responsive UI",
    "All code version-controlled on GitHub",
], Inches(0.55), Inches(2.0), Inches(5.5), Inches(4.5), size=Pt(13), color=OFF_WHITE)

card(s, Inches(6.65), Inches(1.35), Inches(6.3), Inches(5.6))
add_text(s, "🚀  Future Scope",
         Inches(6.85), Inches(1.5), Inches(5.9), Inches(0.45),
         size=Pt(15), bold=True, color=RED_NEON)
bullet_list(s, [
    "Add user login & prediction history",
    "Integrate ECG image analysis (CNN model)",
    "Add more ML models + ensemble voting",
    "Export PDF health report for patients",
    "Doctor dashboard with patient management",
    "Multi-language support",
    "Mobile app (React Native / Flutter)",
    "Integration with wearable health devices",
], Inches(6.85), Inches(2.0), Inches(5.9), Inches(4.5), size=Pt(13), color=OFF_WHITE)

# Bottom thank you
add_rect(s, 0, H-Inches(0.5), W, Inches(0.5), fill=RGBColor(0x1A,0x04,0x04))
add_text(s, "❤️  Thank You  |  Heart Disease Prediction System  |  Built with Python & Scikit-learn",
         Inches(0.5), H-Inches(0.46), Inches(12.3), Inches(0.42),
         size=Pt(12), color=RED_NEON, align=PP_ALIGN.CENTER)


# ── Save ───────────────────────────────────────────────────────
out = "Heart_Disease_Prediction_PPT.pptx"
prs.save(out)
print(f"Saved: {out}")
